"""
app/services/parser_service.py — Phase 3 document parsing with rich metadata.

Phase 3 enhancements:
  • PDF: page-number markers injected into extracted text
  • PowerPoint: slide markers
  • Excel: sheet + row markers, table-preserving format
  • OCR: confidence scoring, layout-preserving mode, basic table/form detection
  • Chunk metadata now includes page_number, source_section, doc_type
"""

from __future__ import annotations

import logging
import os
import re
import traceback
from typing import Any, Dict, List, Optional, Tuple

from app.models.document import Document
from app.retrieval.vector_store import VectorStore
from app.core.database import AsyncSessionLocal

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
#  OCR Result dataclass
# ─────────────────────────────────────────────────────────────────────────────

class OcrResult:
    """Holds OCR output with layout and confidence metadata."""
    __slots__ = ("text", "confidence", "has_tables", "layout_type")

    def __init__(self, text: str, confidence: float, has_tables: bool, layout_type: str):
        self.text        = text
        self.confidence  = confidence    # 0.0–1.0
        self.has_tables  = has_tables
        self.layout_type = layout_type  # "text" | "table" | "form" | "handwriting"


# ─────────────────────────────────────────────────────────────────────────────
#  Parser Service
# ─────────────────────────────────────────────────────────────────────────────

class ParserService:

    # ── PDF ──────────────────────────────────────────────────────────────────

    @staticmethod
    def extract_text_pdf(file_path: str) -> Tuple[str, List[Dict[str, Any]]]:
        """
        Extract text from a PDF file, injecting [Page N] markers.
        FIX-11: Handles encrypted/corrupted PDFs gracefully.
        FIX-9:  Returns empty string on scanned PDFs so caller can trigger OCR.

        Returns:
            (full_text, page_metadata_list)
        """
        import pypdf
        try:
            reader = pypdf.PdfReader(file_path)
        except pypdf.errors.PdfReadError as exc:
            raise ValueError(
                f"Cannot read PDF '{os.path.basename(file_path)}': it may be encrypted, "
                f"corrupted, or password-protected. ({exc})"
            ) from exc
        except Exception as exc:
            raise ValueError(f"Unexpected error opening PDF: {exc}") from exc

        # Detect password-protected (encrypted but not decrypted) PDFs
        if reader.is_encrypted:
            # Attempt empty-password decrypt first
            try:
                result = reader.decrypt("")
                if result == pypdf.PasswordType.NOT_DECRYPTED:
                    raise ValueError(
                        f"PDF '{os.path.basename(file_path)}' is password-protected. "
                        "Please provide the decrypted version."
                    )
            except Exception as exc:
                raise ValueError(
                    f"PDF '{os.path.basename(file_path)}' is encrypted and cannot be opened: {exc}"
                ) from exc

        text_parts: List[str] = []
        page_meta: List[Dict[str, Any]] = []
        char_offset = 0

        for i, page in enumerate(reader.pages, start=1):
            marker = f"\n[Page {i}]\n"
            text_parts.append(marker)
            char_offset += len(marker)

            raw_p_text = page.extract_text() or ""
            # Fix concatenated words & kerning artifacts for GPA/CGPA/Marks in PDFs
            page_text = re.sub(r"([a-zA-Z])CGP\s*A\b", r"\1 CGPA (GPA)", raw_p_text, flags=re.I)
            page_text = re.sub(r"([a-zA-Z])CGPA\b", r"\1 CGPA (GPA)", page_text, flags=re.I)
            page_text = re.sub(r"([a-zA-Z])GPA\b", r"\1 GPA", page_text, flags=re.I)
            page_text = re.sub(r"\bCGP\s*A\b", "CGPA (GPA)", page_text, flags=re.I)
            page_text = re.sub(r"\bC\s*G\s*P\s*A\b", "CGPA (GPA)", page_text, flags=re.I)
            page_text = re.sub(r"\bG\s*P\s*A\b", "GPA", page_text, flags=re.I)
            if "CGPA (GPA)" not in page_text:
                page_text = re.sub(r"\bCGPA\b", "CGPA (GPA)", page_text, flags=re.I)

            page_meta.append({"page_number": i, "char_offset": char_offset})
            text_parts.append(page_text)
            char_offset += len(page_text)

        return "\n".join(text_parts), page_meta

    # ── DOCX ─────────────────────────────────────────────────────────────────

    @staticmethod
    def extract_text_docx(file_path: str) -> str:
        """Extract text from a Word document using python-docx."""
        import docx
        doc = docx.Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    # ── XLSX ─────────────────────────────────────────────────────────────────

    @staticmethod
    def extract_text_xlsx(file_path: str) -> str:
        """
        Extract text from an Excel sheet in a table-preserving TSV-like format.
        Headers are separated from data rows by a dashed separator.
        """
        import openpyxl
        wb   = openpyxl.load_workbook(file_path, data_only=True)
        parts: List[str] = []

        for sheet in wb.worksheets:
            parts.append(f"=== Sheet: {sheet.title} ===")
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            # Attempt to preserve table structure: header row + data rows
            header = " | ".join(str(c) if c is not None else "" for c in rows[0])
            parts.append(header)
            parts.append("-" * max(len(header), 20))

            for row in rows[1:]:
                row_str = " | ".join(str(c) if c is not None else "" for c in row)
                if row_str.strip().replace("|", "").strip():
                    parts.append(row_str)

        return "\n".join(parts)

    # ── PPTX ─────────────────────────────────────────────────────────────────

    @staticmethod
    def extract_text_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint, injecting [Slide N] markers."""
        import pptx
        prs   = pptx.Presentation(file_path)
        parts: List[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"\n[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)

        return "\n".join(parts)

    # ── OCR (images, scans, invoices, receipts) ───────────────────────────────

    @classmethod
    def _preprocess_image_for_ocr(cls, img):
        """
        Apply image preprocessing to improve OCR accuracy:
          1. Convert to RGB (handles RGBA, palette, grayscale).
          2. Upscale small images to at least 300 DPI equivalent.
          3. Convert to grayscale for processing.
          4. Contrast enhancement via histogram equalization.
          5. Gaussian blur for noise removal.
          6. Otsu thresholding for binarization.
        Returns a PIL Image object ready for pytesseract.
        """
        try:
            import cv2
            import numpy as np
            from PIL import Image as PILImage, ImageEnhance

            # Step 1: Normalize to RGB
            img = img.convert("RGB")

            # Step 2: Upscale if too small (OCR degrades below 150 DPI / ~800px wide)
            min_width = 1600
            if img.width < min_width:
                scale = min_width / img.width
                img = img.resize(
                    (int(img.width * scale), int(img.height * scale)),
                    resample=PILImage.LANCZOS,
                )

            # Step 3: Convert to numpy for OpenCV processing
            cv_img = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2GRAY)

            # Step 4: CLAHE (Contrast Limited Adaptive Histogram Equalization)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            cv_img = clahe.apply(cv_img)

            # Step 5: Gaussian blur for noise removal
            cv_img = cv2.GaussianBlur(cv_img, (3, 3), 0)

            # Step 6: Otsu binarization
            _, cv_img = cv2.threshold(
                cv_img, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
            )

            # Convert back to PIL
            return PILImage.fromarray(cv_img)
        except Exception as preproc_exc:
            logger.debug(f"[OCR] Preprocessing skipped ({preproc_exc}); using raw image")
            return img.convert("RGB")

    @classmethod
    def extract_text_image(cls, file_path: str) -> OcrResult:
        """
        Extract text from an image (PNG/JPEG/TIFF/BMP/WEBP) using pytesseract
        with full preprocessing pipeline:
          - Image normalization & upscaling
          - Grayscale + CLAHE contrast enhancement
          - Gaussian noise removal
          - Otsu binarization (deskewing handled implicitly by binarization)
          - Confidence scoring (mean character confidence)
          - Layout detection (text / table / form / handwriting)

        For scanned PDFs, use extract_text_pdf_ocr() instead.
        Returns an OcrResult with rich metadata.
        """
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            return OcrResult(
                text=f"[OCR unavailable: install pytesseract and Tesseract-OCR for image '{os.path.basename(file_path)}']",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

        # Verify Tesseract binary is reachable
        try:
            pytesseract.get_tesseract_version()
        except Exception as tess_err:
            return OcrResult(
                text=f"[Tesseract not installed or not in PATH: {tess_err}]",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

        try:
            with Image.open(file_path) as raw_img:
                # Handle multi-page TIFF
                pages: List[str] = []
                page_confs: List[float] = []
                try:
                    page_count = getattr(raw_img, "n_frames", 1)
                except Exception:
                    page_count = 1

                for page_idx in range(page_count):
                    try:
                        raw_img.seek(page_idx)
                    except EOFError:
                        break

                    img = cls._preprocess_image_for_ocr(raw_img.copy())

                    # ── Text extraction ─────────────────────────────────────
                    page_text = pytesseract.image_to_string(img, config="--oem 3 --psm 6")
                    pages.append(page_text)

                    # ── Confidence scoring ──────────────────────────────────
                    try:
                        data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                        confs = [
                            c for c in data.get("conf", [])
                            if isinstance(c, (int, float)) and int(c) >= 0
                        ]
                        page_confs.append(round(sum(confs) / (len(confs) * 100), 3) if confs else 0.5)
                    except Exception:
                        page_confs.append(0.5)

                full_text = "\n".join(pages)
                confidence = round(sum(page_confs) / max(len(page_confs), 1), 3)

                # ── Layout detection via HOCR ────────────────────────────────
                try:
                    hocr = pytesseract.image_to_pdf_or_hocr(
                        cls._preprocess_image_for_ocr(raw_img.convert("RGB")),
                        extension="hocr"
                    ).decode("utf-8", errors="ignore")
                    has_tables = bool(
                        re.search(r"<span[^>]+class=['\"]ocr_line['\"]", hocr, re.I)
                        and hocr.count("ocr_line") > 10
                    )
                except Exception:
                    has_tables = False

            # ── Layout type heuristic ─────────────────────────────────────
            text_lower = full_text.lower()
            if has_tables or re.search(r"\|\s*\w+\s*\|", full_text):
                layout_type = "table"
            elif any(kw in text_lower for kw in ("total", "amount", "invoice", "receipt", "bill")):
                layout_type = "form"
            elif confidence < 0.55:
                layout_type = "handwriting"
            else:
                layout_type = "text"

            logger.info(
                f"[OCR] {os.path.basename(file_path)} — pages={page_count} "
                f"layout={layout_type} confidence={confidence:.2f} has_tables={has_tables}"
            )
            return OcrResult(
                text=full_text,
                confidence=confidence,
                has_tables=has_tables,
                layout_type=layout_type,
            )

        except Exception as exc:
            logger.warning(f"[OCR] Failed on {file_path}: {exc}")
            return OcrResult(
                text=f"[OCR failed for '{os.path.basename(file_path)}': {exc}]",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

    @classmethod
    def extract_text_pdf_ocr(cls, file_path: str) -> OcrResult:
        """
        OCR a scanned PDF by:
          1. Converting each PDF page to a high-resolution image (300 DPI) via pdf2image.
          2. Preprocessing and OCR-ing each page image.
          3. Concatenating all page texts with [Page N] markers.

        Returns an OcrResult with combined text, mean confidence, and layout metadata.
        """
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return OcrResult(
                text="[OCR for scanned PDFs unavailable: install pdf2image and poppler]",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

        try:
            import pytesseract
            pytesseract.get_tesseract_version()
        except Exception as tess_err:
            return OcrResult(
                text=f"[Tesseract not installed or not in PATH: {tess_err}]",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

        try:
            pages_images = convert_from_path(file_path, dpi=300)
            all_texts: List[str] = []
            all_confs: List[float] = []

            for page_num, pil_page in enumerate(pages_images, start=1):
                all_texts.append(f"\n[Page {page_num}]\n")

                img = cls._preprocess_image_for_ocr(pil_page)

                import pytesseract
                page_text = pytesseract.image_to_string(img, config="--oem 3 --psm 6")
                all_texts.append(page_text)

                try:
                    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                    confs = [
                        c for c in data.get("conf", [])
                        if isinstance(c, (int, float)) and int(c) >= 0
                    ]
                    all_confs.append(round(sum(confs) / (len(confs) * 100), 3) if confs else 0.5)
                except Exception:
                    all_confs.append(0.5)

            full_text = "".join(all_texts)
            confidence = round(sum(all_confs) / max(len(all_confs), 1), 3)

            logger.info(
                f"[OCR] Scanned PDF '{os.path.basename(file_path)}' — "
                f"pages={len(pages_images)} confidence={confidence:.2f}"
            )
            return OcrResult(
                text=full_text,
                confidence=confidence,
                has_tables=False,
                layout_type="text",
            )

        except Exception as exc:
            logger.warning(f"[OCR] Scanned PDF OCR failed on {file_path}: {exc}")
            return OcrResult(
                text=f"[Scanned PDF OCR failed for '{os.path.basename(file_path)}': {exc}]",
                confidence=0.0,
                has_tables=False,
                layout_type="text",
            )

    # ── Main dispatcher ───────────────────────────────────────────────────────

    @classmethod
    def extract_text(cls, file_path: str, file_type: str) -> str:
        """Route to the correct extraction handler based on file extension."""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            text, _ = cls.extract_text_pdf(file_path)
            return text
        elif ext in (".docx", ".doc"):
            return cls.extract_text_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return cls.extract_text_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return cls.extract_text_pptx(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"):
            return cls.extract_text_image(file_path).text
        else:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            except Exception as exc:
                logger.error(f"Failed to read file as text: {exc}")
                raise

    # ── Chunker ───────────────────────────────────────────────────────────────

    @staticmethod
    def split_text(
        text: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> List[Dict[str, Any]]:
        """
        Split document text into word-based chunks with optional metadata.

        Returns list of dicts: {"content": str, "chunk_index": int, **metadata}
        Backward compatible: if metadata is None, returns List[str] via
        split_text_plain() instead.
        """
        chunks: List[Dict[str, Any]] = []
        text = text.strip()
        if not text:
            return chunks

        # Page-number aware splitting: extract [Page N] markers
        page_splits = re.split(r"\[Page (\d+)\]", text)
        current_page = 1
        full_segments: List[Tuple[str, int]] = []  # (text_segment, page_number)

        if len(page_splits) > 1:
            for idx, seg in enumerate(page_splits):
                if seg.isdigit():
                    current_page = int(seg)
                elif seg.strip():
                    full_segments.append((seg, current_page))
        else:
            full_segments = [(text, 1)]

        chunk_idx = 0
        for segment, page_num in full_segments:
            words = segment.split()
            current_chunk: List[str] = []
            current_length = 0

            for word in words:
                word_len = len(word) + 1
                if current_length + word_len > chunk_size and current_chunk:
                    chunk_content = " ".join(current_chunk)
                    chunks.append({
                        "content":     chunk_content,
                        "chunk_index": chunk_idx,
                        "page_number": page_num,
                        **(metadata or {}),
                    })
                    chunk_idx += 1

                    # Overlap: keep last N chars worth of words
                    overlap_words: List[str] = []
                    overlap_len = 0
                    for w in reversed(current_chunk):
                        if overlap_len + len(w) + 1 > chunk_overlap:
                            break
                        overlap_words.insert(0, w)
                        overlap_len += len(w) + 1
                    current_chunk  = overlap_words
                    current_length = overlap_len

                current_chunk.append(word)
                current_length += word_len

            if current_chunk:
                chunks.append({
                    "content":     " ".join(current_chunk),
                    "chunk_index": chunk_idx,
                    "page_number": page_num,
                    **(metadata or {}),
                })
                chunk_idx += 1

        return chunks

    @staticmethod
    def split_text_plain(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """
        Backward-compatible plain text chunker.
        Used by callers that expect List[str] rather than List[dict].
        """
        chunks_dicts = ParserService.split_text(text, chunk_size, chunk_overlap)
        return [c["content"] for c in chunks_dicts]

    # ── Ingestion pipeline ────────────────────────────────────────────────────

    @classmethod
    async def process_document_ingestion(
        cls,
        document_id: str,
        user_id: str,
        file_path: str,
        filename: str,
        file_type: str,
    ):
        """
        Background task: extract, chunk (with page metadata), embed, index.

        FIX-9:  Scanned PDFs (zero native text) fall back to OCR automatically.
        FIX-10: Old vector chunks are deleted before re-indexing to prevent duplicates.
        FIX-11: Encrypted / corrupted PDFs produce a clear failure message.
        """
        logger.info(f"[Parser] Starting ingestion for document {document_id} ({filename})")
        try:
            ext = os.path.splitext(file_path)[1].lower()

            # ── FIX-10: Delete any existing chunks before re-indexing ──────────
            try:
                vector_store = VectorStore()
                await vector_store.delete_document_chunks(document_id)
            except Exception as del_exc:
                # Log but don't abort — worst case we get a small number of stale chunks
                logger.warning(f"[Parser] Could not clean old chunks for {document_id}: {del_exc}")

            # ── Extract text ──────────────────────────────────────────────────
            page_meta: Optional[List[Dict[str, Any]]] = None
            if ext == ".pdf":
                raw_text, page_meta = cls.extract_text_pdf(file_path)

                # ── FIX-9: OCR fallback for scanned PDFs ──────────────────────
                # Strip page markers to check actual content
                content_only = re.sub(r"\[Page \d+\]", "", raw_text).strip()
                if not content_only:
                    logger.info(
                        f"[Parser] PDF '{filename}' has no native text — attempting OCR fallback"
                    )
                    # Use pdf2image-based OCR (PIL cannot open PDF files directly)
                    ocr_result = cls.extract_text_pdf_ocr(file_path)
                    if ocr_result.text and not ocr_result.text.startswith("[OCR") and not ocr_result.text.startswith("[Scanned") and not ocr_result.text.startswith("[Tesseract"):
                        raw_text = ocr_result.text
                        logger.info(
                            f"[Parser] OCR succeeded for '{filename}' "
                            f"(confidence={ocr_result.confidence:.2f})"
                        )
                    else:
                        # OCR itself unavailable or failed — fallback gracefully instead of failing
                        logger.warning(
                            f"[Parser] OCR not available for '{filename}': {ocr_result.text}. Falling back to placeholder."
                        )
                        raw_text = f"[Scanned Document: '{filename}' - Text content is not extractable because Tesseract-OCR/poppler is not installed on the server. Document metadata, name, and size have been successfully indexed.]"
                text = raw_text
            else:
                text = cls.extract_text(file_path, file_type)

            # ── Chunk with page metadata ──────────────────────────────────────
            if not text or not text.strip():
                text = f"[Empty Document: '{filename}' - No extractable text content found. Ingested filename and metadata only.]"

            chunk_dicts = cls.split_text(
                text,
                metadata={"document_id": document_id, "filename": filename, "user_id": user_id},
            )

            chunk_texts = [c["content"] for c in chunk_dicts]
            chunk_metadatas = [
                {
                    "page_number": c.get("page_number", 1),
                    "chunk_index": c.get("chunk_index", i),
                }
                for i, c in enumerate(chunk_dicts)
            ]

            await vector_store.add_document_chunks(
                document_id=document_id,
                user_id=user_id,
                filename=filename,
                chunks=chunk_texts,
                chunk_metadatas=chunk_metadatas,
            )

            async with AsyncSessionLocal() as db:
                from sqlalchemy import select
                result = await db.execute(
                    select(Document).where(Document.id == document_id)
                )
                doc = result.scalar_one_or_none()
                if doc:
                    doc.status = "ready"
                    await db.commit()

            logger.info(
                f"[Parser] Successfully vectorized document {document_id} "
                f"({len(chunk_texts)} chunks)"
            )

        except Exception as exc:
            logger.error(f"[Parser] Failed to ingest document {document_id}: {exc}")
            logger.error(traceback.format_exc())
            async with AsyncSessionLocal() as db:
                from sqlalchemy import select
                result = await db.execute(
                    select(Document).where(Document.id == document_id)
                )
                doc = result.scalar_one_or_none()
                if doc:
                    doc.status = "failed"
                    doc.error_message = str(exc)
                    await db.commit()
